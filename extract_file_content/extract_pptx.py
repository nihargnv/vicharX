from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

if __name__ == "__main__":
    # Test
    pptx_file = "D:\\test\\Unit_2_Chapter_2_CPU_Scheduling_2.pptx"
    extracted_text = extract_text_from_pptx(pptx_file)
    print(extracted_text)
